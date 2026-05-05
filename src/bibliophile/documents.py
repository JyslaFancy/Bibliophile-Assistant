"""
Document processing utilities for various file formats.
"""

import os
import re
from pathlib import Path
from typing import List, Dict, Any, Optional
from rich.console import Console

console = Console()


# Supported file extensions and their handlers
SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "doc",
    ".md": "markdown",
    ".txt": "text",
    ".xlsx": "excel",
    ".xls": "excel",
    ".pptx": "powerpoint",
    ".ppt": "powerpoint",
}


def get_file_type(file_path: str) -> Optional[str]:
    """Determine the file type based on extension."""
    ext = os.path.splitext(file_path)[1].lower()
    return SUPPORTED_EXTENSIONS.get(ext)


class DocumentProcessor:
    """Processes various document formats and extracts text."""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize the document processor.
        
        Args:
            chunk_size: Size of text chunks to create (in characters)
            chunk_overlap: Overlap between chunks (in characters)
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self._ensure_dependencies()
    
    def _ensure_dependencies(self) -> None:
        """Check and warn about missing dependencies."""
        try:
            import pypdf
        except ImportError:
            console.print("[yellow]Warning: pypdf not installed. PDF support disabled.[/yellow]")
        
        try:
            import docx
        except ImportError:
            console.print("[yellow]Warning: python-docx not installed. Word support disabled.[/yellow]")
        
        try:
            import openpyxl
        except ImportError:
            console.print("[yellow]Warning: openpyxl not installed. Excel support disabled.[/yellow]")
        
        try:
            import pptx
        except ImportError:
            console.print("[yellow]Warning: python-pptx not installed. PowerPoint support disabled.[/yellow]")
    
    def scan_folder(self, folder_path: str, extensions: List[str] = None) -> List[str]:
        """
        Scan a folder for supported documents.
        
        Args:
            folder_path: Path to the folder to scan
            extensions: List of extensions to look for (defaults to all supported)
            
        Returns:
            List of file paths
        """
        if extensions is None:
            extensions = list(SUPPORTED_EXTENSIONS.keys())
        
        files = []
        for root, dirs, filenames in os.walk(folder_path):
            for filename in filenames:
                ext = os.path.splitext(filename)[1].lower()
                if ext in extensions:
                    files.append(os.path.join(root, filename))
        
        return files
    
    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process a single file and extract text chunks.
        
        Args:
            file_path: Path to the file to process
            
        Returns:
            List of document chunks with metadata
        """
        file_type = get_file_type(file_path)
        if file_type is None:
            raise ValueError(f"Unsupported file type: {file_path}")
        
        try:
            text = self._extract_text(file_path, file_type)
        except Exception as e:
            raise ValueError(f"Failed to extract text from {file_path}: {e}")
        
        if not text or text.isspace():
            return []
        
        chunks = self._chunk_text(text, file_path, file_type)
        return chunks
    
    def _extract_text(self, file_path: str, file_type: str) -> str:
        """Extract text from a file based on its type."""
        if file_type == "pdf":
            return self._extract_pdf(file_path)
        elif file_type == "docx":
            return self._extract_docx(file_path)
        elif file_type == "doc":
            return self._extract_doc(file_path)
        elif file_type == "markdown":
            return self._extract_markdown(file_path)
        elif file_type == "text":
            return self._extract_text_file(file_path)
        elif file_type == "excel":
            return self._extract_excel(file_path)
        elif file_type == "powerpoint":
            return self._extract_powerpoint(file_path)
        else:
            raise ValueError(f"Unknown file type: {file_type}")
    
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from a PDF file."""
        try:
            import pypdf
            
            text = ""
            with open(file_path, 'rb') as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text
        except ImportError:
            raise ImportError("pypdf is required for PDF support. Install with: pip install pypdf")
        except Exception as e:
            raise ValueError(f"Error reading PDF: {e}")
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from a Word DOCX file."""
        try:
            import docx
            
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text
        except ImportError:
            raise ImportError("python-docx is required for DOCX support. Install with: pip install python-docx")
        except Exception as e:
            raise ValueError(f"Error reading DOCX: {e}")
    
    def _extract_doc(self, file_path: str) -> str:
        """Extract text from a Word DOC file (legacy)."""
        try:
            # Try using textract or antiword if available
            # For now, we'll use a simple approach
            import subprocess
            result = subprocess.run(
                ["antiword", file_path],
                capture_output=True,
                text=True
            )
            if result.returncode == 0:
                return result.stdout
            
            # Fallback: try catdoc
            result = subprocess.run(
                ["catdoc", file_path],
                capture_output=True,
                text=True
            )
            if result.returncode == 0:
                return result.stdout
            
            raise ValueError("No DOC reader available. Install antiword or catdoc.")
        except FileNotFoundError:
            raise ValueError("No DOC reader available. Install antiword or catdoc.")
        except Exception as e:
            raise ValueError(f"Error reading DOC: {e}")
    
    def _extract_markdown(self, file_path: str) -> str:
        """Extract text from a Markdown file."""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _extract_text_file(self, file_path: str) -> str:
        """Extract text from a plain text file."""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _extract_excel(self, file_path: str) -> str:
        """Extract text from an Excel file."""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(file_path)
            text = ""
            
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text += f"\n=== Sheet: {sheet} ===\n"
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text += row_text + "\n"
            
            return text
        except ImportError:
            raise ImportError("openpyxl is required for Excel support. Install with: pip install openpyxl")
        except Exception as e:
            raise ValueError(f"Error reading Excel: {e}")
    
    def _extract_powerpoint(self, file_path: str) -> str:
        """Extract text from a PowerPoint file."""
        try:
            import pptx
            
            prs = pptx.Presentation(file_path)
            text = ""
            
            for i, slide in enumerate(prs.slides):
                text += f"\n=== Slide {i + 1} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text += shape.text + "\n"
            
            return text
        except ImportError:
            raise ImportError("python-pptx is required for PowerPoint support. Install with: pip install python-pptx")
        except Exception as e:
            raise ValueError(f"Error reading PowerPoint: {e}")
    
    def _chunk_text(self, text: str, file_path: str, file_type: str) -> List[Dict[str, Any]]:
        """
        Split text into chunks with metadata.
        
        Args:
            text: The text to chunk
            file_path: Original file path
            file_type: Type of the file
            
        Returns:
            List of chunks with metadata
        """
        # Normalize whitespace
        text = re.sub(r'\s+', ' ', text)
        
        chunks = []
        start = 0
        length = len(text)
        
        while start < length:
            end = min(start + self.chunk_size, length)
            
            # If we're not at the end, try to find a better break point
            if end < length:
                # Look for the last sentence boundary within the overlap
                look_back = min(self.chunk_overlap, end - start)
                for i in range(look_back, 0, -1):
                    if text[start + end - start - i] in ['.', '!', '?', '\n']:
                        end = start + end - start - i + 1
                        break
            
            chunk = text[start:end]
            
            chunks.append({
                "content": chunk,
                "metadata": {
                    "source": file_path,
                    "file_type": file_type,
                    "chunk_index": len(chunks),
                    "start_char": start,
                    "end_char": end
                }
            })
            
            start = end - self.chunk_overlap if end < length else end
        
        return chunks
    
    def process_folder(self, folder_path: str, extensions: List[str] = None) -> List[Dict[str, Any]]:
        """
        Process all documents in a folder.
        
        Args:
            folder_path: Path to the folder
            extensions: List of extensions to process
            
        Returns:
            List of all document chunks from all files
        """
        files = self.scan_folder(folder_path, extensions)
        all_chunks = []
        
        for file_path in files:
            try:
                chunks = self.process_file(file_path)
                all_chunks.extend(chunks)
                console.print(f"[green]Processed: {file_path} ({len(chunks)} chunks)[/green]")
            except Exception as e:
                console.print(f"[red]Error processing {file_path}: {e}[/red]")
        
        return all_chunks
